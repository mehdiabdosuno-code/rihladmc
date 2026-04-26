"""Export engine — PDF (ReportLab), PPTX (python-pptx), XLSX (openpyxl), CSV.
Brand palette from STOURS Studio Guidelines v1.0:
  Bordeaux  #A8371D · Royal Blue #1628A9 · Ink #141414 · Warm Yellow #FFFEE9
"""

import os
import uuid
from pathlib import Path
from datetime import datetime
from typing import Any

from fastapi import HTTPException

from app.modules.reports.schemas import ExportRequest

EXPORT_DIR = Path("exports")
EXPORT_DIR.mkdir(exist_ok=True)

# ── Brand colours (RGB tuples) ────────────────────────────────────
BORDEAUX = (168, 55,  29)   # #A8371D  10 % — CTA / accent
ROYAL    = (22,  40,  169)  # #1628A9   5 % — links / active
INK      = (20,  20,  20)   # #141414  30 % — text / structure
WARM     = (255, 254, 233)  # #FFFEE9  13 % — warm surfaces
WHITE    = (255, 255, 255)  # #FFFFFF  42 % — main bg / print


def _rgb(t: tuple) -> Any:
    """Lazy import helper — returns reportlab Color."""
    from reportlab.lib import colors as rl_colors
    return rl_colors.Color(t[0] / 255, t[1] / 255, t[2] / 255)


# ─────────────────────────────────────────────────────────────────
# PDF via ReportLab
# ─────────────────────────────────────────────────────────────────
def export_pdf(req: ExportRequest) -> Path:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm
        from reportlab.platypus import (
            SimpleDocTemplate, Table, TableStyle,
            Paragraph, Spacer, HRFlowable,
        )
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors as rl_colors
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
    except ImportError:
        raise HTTPException(500, "reportlab non installé — pip install reportlab")

    path = EXPORT_DIR / f"report_{uuid.uuid4().hex[:8]}.pdf"
    doc  = SimpleDocTemplate(
        str(path), pagesize=A4,
        leftMargin=2*cm, rightMargin=2*cm,
        topMargin=2*cm,  bottomMargin=2*cm,
    )

    bord_c = _rgb(BORDEAUX)
    royal_c= _rgb(ROYAL)
    ink_c  = _rgb(INK)
    warm_c = _rgb(WARM)

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "stours_title", parent=styles["Title"],
        fontSize=22, textColor=bord_c,
        fontName="Helvetica-Bold", spaceAfter=4,
    )
    sub_style = ParagraphStyle(
        "stours_sub", parent=styles["Normal"],
        fontSize=10, textColor=rl_colors.Color(0.42, 0.42, 0.42),
        spaceAfter=16,
    )
    body_style = ParagraphStyle(
        "stours_body", parent=styles["Normal"],
        fontSize=10, spaceAfter=6,
    )

    story: list = []

    # Header band
    story.append(Paragraph(req.report_name, title_style))
    story.append(Paragraph(
        req.subtitle or f"Généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')} · S'TOURS Studio",
        sub_style,
    ))
    story.append(HRFlowable(width="100%", thickness=1.5, color=bord_c, spaceAfter=10))
    story.append(Spacer(1, 0.4*cm))

    if req.data and req.fields:
        cols = [f.name for f in req.fields[:9]]
        header_row = cols
        data_rows  = [
            [str(row.get(c, "")) for c in cols]
            for row in req.data[:200]
        ]

        table_data = [header_row] + data_rows
        col_w = (A4[0] - 4*cm) / max(len(cols), 1)

        tbl = Table(table_data, colWidths=[col_w]*len(cols), repeatRows=1)
        tbl.setStyle(TableStyle([
            # Header
            ("BACKGROUND",    (0, 0), (-1,  0), bord_c),
            ("TEXTCOLOR",     (0, 0), (-1,  0), rl_colors.white),
            ("FONTNAME",      (0, 0), (-1,  0), "Helvetica-Bold"),
            ("FONTSIZE",      (0, 0), (-1, -1), 8),
            ("TOPPADDING",    (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            # Rows
            ("ROWBACKGROUNDS", (0, 1), (-1, -1),
             [rl_colors.white, rl_colors.Color(0.98, 0.98, 0.96)]),
            ("GRID",          (0, 0), (-1, -1), 0.3, rl_colors.lightgrey),
            ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
        ]))
        story.append(tbl)
        story.append(Spacer(1, 0.5*cm))

        # Totals line
        num_fields = [f for f in req.fields if f.type == "num"][:4]
        if num_fields:
            totals = "  |  ".join(
                f"{f.name}: "
                f"{sum(float(r.get(f.name, 0) or 0) for r in req.data):,.0f}"
                for f in num_fields
            )
            story.append(Paragraph(f"Totaux — {totals}", body_style))

    story.append(Spacer(1, 1*cm))
    story.append(Paragraph(
        "S'TOURS Studio · S'TOURS DMC Morocco · Confidentiel",
        sub_style,
    ))

    doc.build(story)
    return path


# ─────────────────────────────────────────────────────────────────
# PPTX via python-pptx — STOURS Studio brand
# ─────────────────────────────────────────────────────────────────
def export_pptx(req: ExportRequest) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise HTTPException(500, "python-pptx non installé — pip install python-pptx")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def rgb(t): return RGBColor(*t)

    def rect(slide, l, t, w, h, fill=None, line_color=None):
        s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        s.line.fill.background()
        if fill:
            s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
        else:
            s.fill.background()
        if line_color:
            s.line.color.rgb = rgb(line_color)
        else:
            s.line.fill.background()
        return s

    def text(slide, txt, l, t, w, h, size=14, bold=False,
             color=INK, align=PP_ALIGN.LEFT):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf  = txb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(txt)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return txb

    # ── Slide 1 · Cover ──────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    rect(s1, 0, 0, 13.33, 7.5, fill=INK)           # full dark bg
    rect(s1, 0, 0, 0.3,   7.5, fill=BORDEAUX)       # left accent bar
    rect(s1, 0.3, 6.2, 13.03, 0.04, fill=BORDEAUX)  # bottom line

    # Logo monogram
    rect(s1, 0.7, 0.55, 0.72, 0.72, fill=BORDEAUX)
    text(s1, "S'T", 0.72, 0.55, 0.72, 0.72, size=16, bold=True, color=WARM,
         align=PP_ALIGN.CENTER)

    text(s1, "S'TOURS Studio", 1.6, 0.55, 8, 0.45, size=13, color=(160, 160, 160))
    text(s1, req.report_name,  0.7, 1.5,  12, 1.6,  size=36, bold=True, color=WARM)
    sub = req.subtitle or f"S'TOURS DMC Morocco · {datetime.now().strftime('%d/%m/%Y')}"
    text(s1, sub, 0.7, 3.25, 10, 0.5, size=14, color=(160, 160, 160))
    text(s1, f"{len(req.data)} enregistrements · {len(req.fields)} indicateurs",
         0.7, 4.0, 8, 0.4, size=11, color=(100, 100, 100))

    # ── Slide 2 · KPIs ───────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    rect(s2, 0, 0, 13.33, 7.5, fill=(250, 250, 248))
    rect(s2, 0, 0, 13.33, 1.05, fill=BORDEAUX)
    rect(s2, 0, 0, 0.18,  1.05, fill=INK)
    text(s2, "Indicateurs clés", 0.35, 0.18, 10, 0.7, size=22, bold=True, color=WARM)

    num_fields = [f for f in req.fields if f.type == "num"][:4]
    kpi_x = [0.4, 3.55, 6.7, 9.85]
    for i, f in enumerate(num_fields):
        vals  = [float(r.get(f.name, 0) or 0) for r in req.data]
        total = sum(vals)
        x     = kpi_x[i]
        rect(s2, x, 1.4, 2.95, 2.1, fill=WHITE, line_color=(220, 220, 215))
        rect(s2, x, 1.4, 2.95, 0.16, fill=BORDEAUX)
        disp  = f"{total/1000:.0f}k" if total >= 10_000 else f"{total:,.0f}"
        text(s2, disp,   x+0.15, 1.72, 2.65, 0.9, size=32, bold=True, color=INK, align=PP_ALIGN.CENTER)
        text(s2, f.label or f.name, x+0.15, 2.7, 2.65, 0.45, size=11,
             color=(120, 120, 120), align=PP_ALIGN.CENTER)

    # ── Slide 3 · Data table ──────────────────────────────────────
    if req.data and req.fields:
        s3   = prs.slides.add_slide(blank)
        rect(s3, 0, 0, 13.33, 7.5, fill=(250, 250, 248))
        rect(s3, 0, 0, 13.33, 1.05, fill=BORDEAUX)
        rect(s3, 0, 0, 0.18,  1.05, fill=INK)
        text(s3, "Données détaillées", 0.35, 0.18, 10, 0.7, size=22, bold=True, color=WARM)

        cols   = req.fields[:7]
        n_cols = len(cols)
        col_w  = 12.3 / max(n_cols, 1)
        sx, sy = 0.5, 1.15
        row_h  = 0.48

        # Header row
        for ci, f in enumerate(cols):
            x = sx + ci * col_w
            rect(s3, x, sy, col_w - 0.04, row_h, fill=INK)
            text(s3, f.label or f.name, x+0.05, sy+0.06, col_w-0.1, row_h-0.1,
                 size=9, bold=True, color=WARM)

        # Data rows (max 10)
        for ri, row in enumerate(req.data[:10]):
            y  = sy + (ri + 1) * row_h
            bg = WHITE if ri % 2 == 0 else (246, 245, 240)
            for ci, f in enumerate(cols):
                x   = sx + ci * col_w
                val = str(row.get(f.name, ""))
                rect(s3, x, y, col_w - 0.04, row_h - 0.03, fill=bg)
                text(s3, val, x+0.05, y+0.07, col_w-0.1, row_h-0.1, size=8, color=INK)

    path = EXPORT_DIR / f"report_{uuid.uuid4().hex[:8]}.pptx"
    prs.save(str(path))
    return path


# ─────────────────────────────────────────────────────────────────
# XLSX via openpyxl
# ─────────────────────────────────────────────────────────────────
def export_xlsx(req: ExportRequest) -> Path:
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.utils import get_column_letter
    except ImportError:
        raise HTTPException(500, "openpyxl non installé — pip install openpyxl")

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Rapport"

    bord_hex = "A8371D"
    ink_hex  = "141414"
    warm_hex = "FFFEE9"
    alt_hex  = "F9F8F3"

    cols = [f.name for f in req.fields]

    # Header
    for ci, col in enumerate(cols, 1):
        cell = ws.cell(row=1, column=ci, value=col)
        cell.font      = Font(bold=True, color="FFFFFF", size=10, name="Inter")
        cell.fill      = PatternFill("solid", fgColor=bord_hex)
        cell.alignment = Alignment(horizontal="center", vertical="center")
        ws.column_dimensions[get_column_letter(ci)].width = max(len(col) + 4, 14)

    ws.row_dimensions[1].height = 26

    # Data
    for ri, row in enumerate(req.data, 2):
        fill_hex = warm_hex if ri % 2 == 0 else alt_hex
        for ci, col in enumerate(cols, 1):
            cell = ws.cell(row=ri, column=ci, value=row.get(col, ""))
            cell.fill      = PatternFill("solid", fgColor=fill_hex)
            cell.alignment = Alignment(vertical="center")
            cell.font      = Font(name="Inter", size=9)

    # Totals row
    total_row = len(req.data) + 2
    for ci, f in enumerate(req.fields, 1):
        if f.type == "num":
            total = sum(float(r.get(f.name, 0) or 0) for r in req.data)
            cell  = ws.cell(row=total_row, column=ci, value=round(total, 2))
        elif ci == 1:
            cell = ws.cell(row=total_row, column=ci, value=f"Total ({len(req.data)})")
        else:
            cell = ws.cell(row=total_row, column=ci, value="")
        cell.font = Font(bold=True, color="FFFFFF", name="Inter", size=9)
        cell.fill = PatternFill("solid", fgColor=ink_hex)

    path = EXPORT_DIR / f"report_{uuid.uuid4().hex[:8]}.xlsx"
    wb.save(str(path))
    return path


# ─────────────────────────────────────────────────────────────────
# CSV
# ─────────────────────────────────────────────────────────────────
def export_csv(req: ExportRequest) -> Path:
    import csv
    path = EXPORT_DIR / f"report_{uuid.uuid4().hex[:8]}.csv"
    cols = [f.name for f in req.fields]
    with open(path, "w", newline="", encoding="utf-8-sig") as fh:
        writer = csv.DictWriter(fh, fieldnames=cols, extrasaction="ignore")
        writer.writeheader()
        writer.writerows(req.data)
    return path
